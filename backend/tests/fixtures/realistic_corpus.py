"""Programmatically generated realistic structural test corpus for Phase 2.

Corpus Version: phase2-structural-corpus-v1
Generates real structural constructs (multi-page, H1/H2 hierarchies, paragraphs, lists, tables)
across PDF, DOCX, PPTX, Markdown, Python code, CSV, JSON, and XLSX.
"""

import csv
import json
import os
from typing import Dict
import docx
import openpyxl
import pptx
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors

CORPUS_VERSION = "phase2-structural-corpus-v1"


def generate_realistic_structural_corpus(target_dir: str) -> Dict[str, str]:
    """
    Creates a full set of realistic structural document fixtures in target_dir.
    Returns a dict mapping format names to their absolute file paths.
    """
    os.makedirs(target_dir, exist_ok=True)
    paths = {}

    # 1. Multi-page Structural PDF
    pdf_path = os.path.join(target_dir, "sample_system_spec.pdf")
    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=letter,
        rightMargin=54,
        leftMargin=54,
        topMargin=54,
        bottomMargin=54,
    )
    styles = getSampleStyleSheet()
    h1_style = ParagraphStyle(
        "H1Style",
        parent=styles["Heading1"],
        fontSize=20,
        leading=24,
        textColor=colors.HexColor("#1e293b"),
        spaceAfter=12,
    )
    h2_style = ParagraphStyle(
        "H2Style",
        parent=styles["Heading2"],
        fontSize=14,
        leading=18,
        textColor=colors.HexColor("#334155"),
        spaceAfter=8,
    )
    body_style = ParagraphStyle(
        "BodyStyle",
        parent=styles["Normal"],
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#475569"),
        spaceAfter=8,
    )

    story = []
    # Page 1
    story.append(Paragraph("System Architecture Specification", h1_style))
    story.append(Paragraph("This document specifies the internal pipeline architecture for the FileMind document intelligence engine.", body_style))
    story.append(Spacer(1, 10))
    story.append(Paragraph("Component Overview", h2_style))
    story.append(Paragraph("The system is divided into modular subsystems:", body_style))
    story.append(Paragraph("• Filesystem Watcher Service", body_style))
    story.append(Paragraph("• Worker Queue and State Persistence", body_style))
    story.append(Paragraph("• Hierarchical Chunker with Provenance Tracking", body_style))
    story.append(Spacer(1, 14))

    # Page 1 Table
    table_data = [
        ["Subsystem", "Latency Target", "Status"],
        ["Watcher Debounce", "500 ms", "Operational"],
        ["Crash Recovery", "< 50 ms", "Operational"],
        ["Parser Dispatch", "< 100 ms", "Operational"],
    ]
    t = Table(table_data, colWidths=[160, 120, 100])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#f1f5f9")),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor("#0f172a")),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor("#cbd5e1")),
    ]))
    story.append(t)
    story.append(PageBreak())

    # Page 2
    story.append(Paragraph("Storage Engine Specifications", h1_style))
    story.append(Paragraph("The persistence tier utilizes SQLite with Write-Ahead Logging (WAL) mode enabled for high-concurrency read operations.", body_style))
    story.append(Spacer(1, 10))
    story.append(Paragraph("Data Integrity Protocols", h2_style))
    story.append(Paragraph("Every chunk persisted carries an immutable provenance record including source file ID, page offsets, section hierarchy, and cryptographic content hashes.", body_style))
    
    doc.build(story)
    paths["PDF"] = pdf_path

    # 2. Structural DOCX
    docx_path = os.path.join(target_dir, "sample_system_spec.docx")
    d_doc = docx.Document()
    d_doc.core_properties.title = "System Architecture Specification"
    d_doc.core_properties.author = "FileMind Architecture Team"

    d_doc.add_heading("System Architecture Specification", level=1)
    d_doc.add_paragraph("This document outlines the core architecture and chunk provenance model for FileMind.")
    
    d_doc.add_heading("Component Overview", level=2)
    d_doc.add_paragraph("Key architectural components include:", style="List Bullet")
    d_doc.add_paragraph("Filesystem Watcher Service", style="List Bullet")
    d_doc.add_paragraph("Worker Pool with Backoff", style="List Bullet")
    d_doc.add_paragraph("Hierarchical Document Chunker", style="List Bullet")

    # DOCX Table
    tbl = d_doc.add_table(rows=1, cols=3)
    hdr_cells = tbl.rows[0].cells
    hdr_cells[0].text = "Component"
    hdr_cells[1].text = "Layer"
    hdr_cells[2].text = "Persistence"
    
    row_data = [
        ("Watcher", "Filesystem", "SQLite file_events"),
        ("WorkerPool", "Engine", "SQLite indexing_jobs"),
        ("ChunkStore", "Intelligence", "SQLite chunks"),
    ]
    for c, l, p in row_data:
        row_cells = tbl.add_row().cells
        row_cells[0].text = c
        row_cells[1].text = l
        row_cells[2].text = p

    d_doc.add_heading("Storage Engine Specifications", level=1)
    d_doc.add_paragraph("The storage engine leverages SQLite with Write-Ahead Logging (WAL) for persistent durability.")
    d_doc.save(docx_path)
    paths["DOCX"] = docx_path

    # 3. Structural PPTX
    pptx_path = os.path.join(target_dir, "sample_presentation.pptx")
    prs = pptx.Presentation()

    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "FileMind Architecture Overview"
    slide1.placeholders[1].text = "Phase 2 Document Intelligence & Provenance Subsystem"

    # Slide 2 (Bullets & Table)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Core Processing Pipeline"
    tf = slide2.placeholders[1].text_frame
    tf.text = "Pipeline Stages:"
    p1 = tf.add_paragraph()
    p1.text = "• Discovery and SHA-256 Verification"
    p1.level = 1
    p2 = tf.add_paragraph()
    p2.text = "• Parser Registry and Layout Extraction"
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "• Hierarchical Chunking and Provenance Builder"
    p3.level = 1

    prs.save(pptx_path)
    paths["PPTX"] = pptx_path

    # 4. Markdown Document
    md_path = os.path.join(target_dir, "sample_architecture.md")
    md_content = """# FileMind Engine Specification

## Overview
FileMind is a local-first desktop intelligence system designed for local filesystem discovery and structural document processing.

## Subsystems
The system consists of the following components:
- **Filesystem Engine**: SQLite WAL persistence and change detection.
- **Document Intelligence**: Format-aware extraction and hierarchical chunking.

### Configuration
```python
def get_config():
    return {
        "integrity_mode": "NORMAL",
        "debounce_ms": 500,
        "max_workers": 4,
    }
```

## Performance Targets
| Subsystem | Metric | Target |
|---|---|---|
| Watcher | Latency | < 600 ms |
| Recovery | Latency | < 50 ms |
| Chunker | Throughput | > 500 chunks/s |

## Conclusion
This concludes the architectural specification.
"""
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(md_content)
    paths["MARKDOWN"] = md_path

    # 5. Python Source Code
    py_path = os.path.join(target_dir, "sample_coordinator.py")
    py_content = '''"""Coordinator module for FileMind background processing."""

import os
import time

class EngineCoordinator:
    """Coordinates filesystem discovery and worker threads."""

    def __init__(self, db_manager):
        self.db = db_manager
        self.is_active = True

    def initialize(self):
        """Initializes database schema migrations and worker pool."""
        return True

    def execute_task(self, task_id: str):
        """Executes a single processing task."""
        time.sleep(0.01)
        return {"status": "COMPLETED", "task_id": task_id}
'''
    with open(py_path, "w", encoding="utf-8") as f:
        f.write(py_content)
    paths["CODE"] = py_path

    # 6. CSV Dataset
    csv_path = os.path.join(target_dir, "sample_metrics.csv")
    with open(csv_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["metric_name", "category", "target_val", "unit"])
        writer.writerow(["discovery_rate", "throughput", "500", "files/s"])
        writer.writerow(["sha256_rate", "throughput", "50", "MB/s"])
        writer.writerow(["watcher_latency", "latency", "600", "ms"])
        writer.writerow(["recovery_time", "latency", "50", "ms"])
    paths["CSV"] = csv_path

    # 7. JSON Dataset
    json_path = os.path.join(target_dir, "sample_dataset.json")
    json_content = [
        {"id": "doc_001", "name": "Architecture Spec", "pages": 12, "format": "PDF"},
        {"id": "doc_002", "name": "Storage Engine", "pages": 8, "format": "DOCX"},
        {"id": "doc_003", "name": "Metrics Report", "pages": 3, "format": "XLSX"},
    ]
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(json_content, f, indent=2)
    paths["JSON"] = json_path

    # 8. XLSX Workbook
    xlsx_path = os.path.join(target_dir, "sample_metrics.xlsx")
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "PerformanceBenchmarks"
    ws1.append(["Benchmark", "Runs", "Median", "Unit"])
    ws1.append(["Discovery Throughput", 5, 761.56, "files/s"])
    ws1.append(["SHA-256 Throughput", 5, 606.74, "MB/s"])
    ws1.append(["Watcher Latency", 5, 556.51, "ms"])

    ws2 = wb.create_sheet(title="SystemCapacity")
    ws2.append(["Resource", "Baseline", "Headroom"])
    ws2.append(["Startup Latency", "3.705 s", "1.295 s"])
    ws2.append(["Process RSS", "27.62 MB", "172.38 MB"])
    wb.save(xlsx_path)
    wb.close()
    paths["XLSX"] = xlsx_path

    return paths
