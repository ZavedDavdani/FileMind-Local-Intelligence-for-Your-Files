"""PPTX presentation parser using python-pptx."""

import os
from typing import List, Optional
import pptx

from app.intelligence.models import (
    Document,
    DocumentElement,
    ElementType,
    TableData,
)
from app.intelligence.parsers.base import (
    BaseParser,
    CorruptedDocumentError,
    DocumentParserError,
)


class PptxParser(BaseParser):
    """
    Parser for Microsoft PowerPoint (.pptx) presentations.
    Extracts slide numbers, slide titles, body shapes, bullet lists, and tables.
    """

    @property
    def parser_name(self) -> str:
        return "pptx-parser"

    @property
    def parser_version(self) -> str:
        return "1.0.0"

    @property
    def supported_mime_types(self) -> List[str]:
        return ["application/vnd.openxmlformats-officedocument.presentationml.presentation"]

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx"]

    def parse(self, file_path: str, file_id: str, mime_type: str = "application/vnd.openxmlformats-officedocument.presentationml.presentation") -> Document:
        filename = os.path.basename(file_path)
        doc_obj = Document(
            file_id=file_id,
            source_path=file_path,
            filename=filename,
            mime_type=mime_type,
            parser_name=self.parser_name,
            parser_version=self.parser_version,
        )

        try:
            prs = pptx.Presentation(file_path)
        except Exception as exc:
            raise CorruptedDocumentError(f"Failed to open PPTX file: {str(exc)}") from exc

        doc_obj.total_pages = len(prs.slides)
        element_idx = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_title = None
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()

            current_slide_heading_id = None
            if slide_title:
                element_idx += 1
                elem_id = f"{file_id}_elem_{element_idx}"
                elem = DocumentElement(
                    element_id=elem_id,
                    element_type=ElementType.HEADING,
                    text=f"Slide {slide_idx}: {slide_title}",
                    page_number=slide_idx,
                    level=1,
                )
                current_slide_heading_id = elem_id
                doc_obj.elements.append(elem)

            for shape in slide.shapes:
                # Skip title shape if already processed
                if shape == slide.shapes.title:
                    continue

                if shape.has_table:
                    # Extract slide table
                    table = shape.table
                    rows_data = []
                    for r in table.rows:
                        rows_data.append([c.text.strip() for c in r.cells])

                    if rows_data:
                        headers = rows_data[0]
                        rows = rows_data[1:] if len(rows_data) > 1 else []
                        t_data = TableData(headers=headers, rows=rows)
                        element_idx += 1
                        elem = DocumentElement(
                            element_id=f"{file_id}_elem_{element_idx}",
                            element_type=ElementType.TABLE,
                            text=t_data.to_markdown(),
                            page_number=slide_idx,
                            table_data=t_data,
                            parent_heading_id=current_slide_heading_id,
                        )
                        doc_obj.elements.append(elem)

                elif shape.has_text_frame:
                    # Extract paragraphs and bullet points
                    for p in shape.text_frame.paragraphs:
                        text = p.text.strip()
                        if not text:
                            continue

                        element_idx += 1
                        elem_id = f"{file_id}_elem_{element_idx}"

                        if p.level > 0 or text.startswith(("- ", "• ", "* ")):
                            elem = DocumentElement(
                                element_id=elem_id,
                                element_type=ElementType.LIST_ITEM,
                                text=text,
                                page_number=slide_idx,
                                parent_heading_id=current_slide_heading_id,
                            )
                        else:
                            elem = DocumentElement(
                                element_id=elem_id,
                                element_type=ElementType.PARAGRAPH,
                                text=text,
                                page_number=slide_idx,
                                parent_heading_id=current_slide_heading_id,
                            )
                        doc_obj.elements.append(elem)

        return doc_obj
